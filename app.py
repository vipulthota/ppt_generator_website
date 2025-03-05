import os
import re
import requests
from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Flask App Initialization
app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# OpenAI API Key
OPENAI_API_KEY = "Your_OpenAI_API_Key"
# Function to generate slide content using OpenAI API
def generate_slide_content(topic, num_slides):
    prompt = (
        f"Create a structured {num_slides}-slide PowerPoint presentation on '{topic}'. "
        "Each slide must follow this format:\n"
        "'Slide X:\nTitle: <Slide Title>\nContent: <Detailed content>'\n"
        "Ensure content is informative and well-structured, with no bullet points."
    )

    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json"
    }

    data = {
        "model": "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": "You are an expert at generating structured PowerPoint presentations."},
            {"role": "user", "content": prompt}
        ]
    }

    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=data)

    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        print(f"Error: {response.json()}")
        return None

# Function to parse slide content
def parse_slide_content(content):
    return re.findall(r"Slide \d+:\s*Title:\s*(.*?)\s*Content:\s*(.*?)\n", content, re.DOTALL)

# Function to create PowerPoint presentation
def create_ppt(topic, slide_content):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = topic

    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_paragraph.font.bold = True
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.name = "Arial"
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.font.color.rgb = RGBColor(0, 51, 102)

    # Generate slides
    slides_data = parse_slide_content(slide_content)

    for title_text, content_text in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.shapes.placeholders[1]

        title.text = title_text
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.name = "Arial"
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        content.text = ""
        for line in content_text.split(". "):
            p = content.text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(18)
            p.font.name = "Arial"
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)

    output_pptx = os.path.join(UPLOAD_FOLDER, f"{topic.replace(' ', '_')}.pptx")
    prs.save(output_pptx)
    return output_pptx

# Route for Home Page
@app.route('/')
def home():
    return render_template("index.html")

# Route to generate and download the PowerPoint
@app.route('/generate', methods=['POST'])
def generate():
    topic = request.form['topic']
    num_slides = int(request.form['num_slides'])
    
    slide_content = generate_slide_content(topic, num_slides)
    
    if slide_content:
        ppt_file = create_ppt(topic, slide_content)
        return send_file(ppt_file, as_attachment=True)

    return "Failed to generate the PowerPoint.", 500

# Run Flask App
if __name__ == '__main__':
    app.run(debug=True)
